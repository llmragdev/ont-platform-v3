from __future__ import annotations

from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm, Pt


REPORT_DIR = Path("E:/ontology_edu/X_ont_std/validation/ont_platform_v4_eval/reports")
OLD_DIR = Path("E:/ai_lab_SIT/reports")
OUT = REPORT_DIR / "AI Lab 2기 중간평가 API 정확도 리뷰.pptx"
TEAM1_WRONG_SOURCE_ADJUSTED = 76.88


COLORS = {
    "ink": RGBColor(30, 41, 59),
    "muted": RGBColor(100, 116, 139),
    "navy": RGBColor(21, 42, 78),
    "blue": RGBColor(37, 99, 235),
    "teal": RGBColor(20, 184, 166),
    "green": RGBColor(22, 163, 74),
    "orange": RGBColor(234, 88, 12),
    "red": RGBColor(220, 38, 38),
    "amber": RGBColor(245, 158, 11),
    "line": RGBColor(203, 213, 225),
    "soft": RGBColor(241, 245, 249),
    "white": RGBColor(255, 255, 255),
}


def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def find_file(folder: Path, keyword: str) -> Path:
    files = [p for p in folder.glob("*.xlsx") if keyword in p.name and not p.name.startswith("~$")]
    if not files:
        raise FileNotFoundError(keyword)
    return files[0]


def read_summary() -> dict:
    old_file = find_file(OLD_DIR, "답변 보정전")
    new_file = find_file(REPORT_DIR, "답변 보정후")
    old_wb = load_workbook(str(old_file), read_only=True, data_only=True)
    new_wb = load_workbook(str(new_file), read_only=True, data_only=True)

    def sheet_by_header(wb, first_header: str, second_header_contains: str):
        for ws in wb.worksheets:
            if ws.cell(1, 1).value == first_header and second_header_contains in str(ws.cell(1, 2).value):
                return ws
        raise RuntimeError(f"sheet not found: {first_header}, {second_header_contains}")

    old_team = sheet_by_header(old_wb, "팀", "정확도")
    old_cat = sheet_by_header(old_wb, "카테고리", "Team0")
    new_team = sheet_by_header(new_wb, "팀", "전체")
    new_cat = sheet_by_header(new_wb, "카테고리", "Team0")
    detail = next(ws for ws in new_wb.worksheets if ws.cell(1, 1).value == "문제ID" and ws.max_column >= 13)

    label_map = {
        "Team0": "0조",
        "Team1": "1조",
        "Team2": "2조",
        "Team4": "4조",
        "Team4 (ont_platform v4)": "4조",
    }

    old_scores = {}
    for row in range(2, old_team.max_row + 1):
        name = str(old_team.cell(row, 1).value)
        old_scores[label_map.get(name, name)] = float(old_team.cell(row, 2).value)

    new_scores = {}
    for row in range(2, new_team.max_row + 1):
        name = str(new_team.cell(row, 1).value).replace(" (ont_platform v4)", "")
        new_scores[label_map.get(name, name)] = float(new_team.cell(row, 2).value)

    categories = []
    for row in range(2, new_cat.max_row + 1):
        category = new_cat.cell(row, 1).value
        if category == "전체 평균":
            continue
        categories.append(
            {
                "category": str(category),
                "Team0": float(new_cat.cell(row, 2).value),
                "Team1": float(new_cat.cell(row, 3).value),
                "Team2": float(new_cat.cell(row, 4).value),
                "Team4": float(new_cat.cell(row, 5).value),
                "avg": float(new_cat.cell(row, 6).value),
            }
        )

    std_s = []
    for row in range(2, detail.max_row + 1):
        qid = detail.cell(row, 1).value
        if qid and str(qid).startswith("STD-S"):
            std_s.append(
                [
                    str(qid),
                    float(detail.cell(row, 5).value),
                    float(detail.cell(row, 7).value),
                    float(detail.cell(row, 9).value),
                    float(detail.cell(row, 11).value),
                ]
            )

    return {
        "old_file": old_file,
        "new_file": new_file,
        "old_scores": old_scores,
        "new_scores": new_scores,
        "categories": categories,
        "std_s": std_s,
    }


def set_font(run, size=18, bold=False, color=None):
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color or COLORS["ink"])
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, 1.05, 0.55, 30.0, 1.05, title, 25, True, COLORS["navy"])
    if subtitle:
        add_text(slide, 1.08, 1.46, 29.5, 0.62, subtitle, 10.5, False, COLORS["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.05), Cm(2.05), Cm(31.0), Cm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_footer(slide, idx):
    add_text(slide, 29.3, 18.45, 2.1, 0.35, f"{idx:02d}", 8, False, COLORS["muted"], PP_ALIGN.RIGHT)
    add_text(slide, 1.05, 18.45, 12.0, 0.35, "AI Lab 2기 중간평가 API 정확도 리뷰", 8, False, COLORS["muted"])


def add_pill(slide, x, y, w, text, fill, color=COLORS["white"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_font(r, size=8.5, bold=True, color=color)
    return shape


def add_rule(slide, x, y, w, color=COLORS["line"]):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_open_metric(slide, x, y, label, value, note, color):
    add_text(slide, x, y, 4.2, 0.45, label, 9, True, COLORS["muted"])
    add_text(slide, x, y + 0.35, 4.4, 0.88, value, 25, True, color)
    add_text(slide, x, y + 1.25, 4.9, 0.55, note, 8.2, False, COLORS["muted"])
    add_rule(slide, x, y + 1.95, 4.7, COLORS["line"])


def add_bullets(slide, x, y, w, h, items, size=12.0, color=COLORS["ink"], gap=0.08):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = f"• {item}"
        set_font(r, size=size, color=color)
    return box


def color_for_score(score):
    if score >= 90:
        return COLORS["green"]
    if score >= 80:
        return COLORS["teal"]
    if score >= 70:
        return COLORS["amber"]
    if score >= 50:
        return COLORS["orange"]
    return COLORS["red"]


def add_heatmap(slide, x, y, categories):
    teams = ["Team0", "Team1", "Team2", "Team4"]
    labels = {"Team0": "0조", "Team1": "1조", "Team2": "2조", "Team4": "4조"}
    col_w = [4.2, 3.2, 3.2, 3.2, 3.2]
    row_h = 1.08
    headers = ["카테고리", *[labels[t] for t in teams]]
    for c, header in enumerate(headers):
        add_text(slide, x + sum(col_w[:c]), y, col_w[c] - 0.1, 0.45, header, 8.8, True, COLORS["muted"], PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
    for r, item in enumerate(categories):
        yy = y + 0.72 + r * row_h
        add_text(slide, x, yy + 0.26, col_w[0] - 0.2, 0.42, item["category"], 10.2, True, COLORS["ink"])
        for c, team in enumerate(teams, start=1):
            score = item[team]
            fill = color_for_score(score)
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x + sum(col_w[:c])), Cm(yy), Cm(col_w[c] - 0.25), Cm(0.82))
            rect.fill.solid()
            rect.fill.fore_color.rgb = fill
            rect.line.fill.background()
            rect.text_frame.clear()
            p = rect.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r0 = p.add_run()
            r0.text = f"{score:.2f}"
            set_font(r0, size=11, bold=True, color=COLORS["white"])


def add_team_chart(slide, x, y, w, h, scores):
    chart_data = ChartData()
    chart_data.categories = list(scores.keys())
    chart_data.add_series("정확도 (%)", list(scores.values()))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Cm(x), Cm(y), Cm(w), Cm(h), chart_data).chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.major_unit = 20
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.plots[0].has_data_labels = True
    labels = chart.plots[0].data_labels
    labels.number_format = "0.0"
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    labels.font.size = Pt(8)
    for i, point in enumerate(chart.series[0].points):
        score = list(scores.values())[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color_for_score(score)
    return chart


def add_change_chart(slide, x, y, w, h, old_scores, new_scores):
    teams = ["0조", "1조", "2조"]
    chart_data = ChartData()
    chart_data.categories = teams
    chart_data.add_series("기존 3팀 평가", [old_scores[t] for t in teams])
    chart_data.add_series("답변 보정후 산정", [new_scores[t] for t in teams])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Cm(x), Cm(y), Cm(w), Cm(h), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.number_format = "0.0"
    chart.plots[0].data_labels.font.size = Pt(8)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(148, 163, 184)
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = COLORS["blue"]
    return chart


def add_simple_table(slide, x, y, widths, rows, header=True, font_size=8.6):
    row_h = 0.78
    for r, row in enumerate(rows):
        yy = y + r * row_h
        for c, text in enumerate(row):
            xx = x + sum(widths[:c])
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(xx), Cm(yy), Cm(widths[c]), Cm(row_h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = COLORS["navy"] if header and r == 0 else (RGBColor(248, 250, 252) if r % 2 else COLORS["white"])
            rect.line.color.rgb = COLORS["line"]
            rect.line.width = Pt(0.5)
            rect.text_frame.clear()
            rect.text_frame.margin_left = Cm(0.12)
            rect.text_frame.margin_right = Cm(0.12)
            rect.text_frame.margin_top = Cm(0.07)
            p = rect.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            rr = p.add_run()
            rr.text = str(text)
            set_font(rr, font_size, bold=(header and r == 0), color=COLORS["white"] if header and r == 0 else COLORS["ink"])


def build_deck():
    data = read_summary()
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    # 1 cover
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_pill(slide, 1.2, 1.15, 4.0, "답변 보정후", COLORS["blue"])
    add_text(slide, 1.18, 3.15, 21.0, 2.6, "AI Lab 2기 중간평가\nAPI 정확도 리뷰", 30, True, COLORS["navy"])
    add_text(slide, 1.22, 6.35, 21.5, 0.8, "예상 답안 보정, 함정 질문 처리, 근거 기반 평가 방식 보완", 13, False, COLORS["muted"])
    add_rule(slide, 1.22, 8.05, 13.0, COLORS["blue"])
    add_text(slide, 1.25, 9.25, 20.0, 1.1, "보정 전 평가는 시스템 간 단순 정확도 비교였다.", 21, True, COLORS["navy"])
    add_text(slide, 1.25, 10.75, 20.5, 2.4, "보정 후 평가는 ‘근거가 없는 질문을 거절할 수 있는가’, ‘잘못된 출처를 차단할 수 있는가’, ‘검색 근거를 설명할 수 있는가’를 함께 보는 평가로 바뀌었다.", 18, True, COLORS["ink"])
    add_text(slide, 23.0, 3.0, 8.2, 6.2, "이번 리뷰는 순위 발표보다 평가 기준의 보정과 서비스화 시 고려해야 할 근거 통제 원칙을 정리하는 데 목적이 있다.", 18, True, COLORS["ink"])
    add_pill(slide, 23.0, 10.0, 3.8, "핵심 관점", COLORS["teal"])
    add_bullets(slide, 23.0, 10.95, 8.8, 3.8, [
        "근거 없는 질문 거절",
        "잘못된 출처 차단",
        "검색 근거 설명 가능성",
    ], 13.5)
    add_text(slide, 1.2, 17.4, 16.0, 0.45, "Source: AI Lab 2기 중간평가 API 정확도 산정(답변 보정전/후).xlsx", 8.2, False, COLORS["muted"])

    # 2 overview
    slide = prs.slides.add_slide(blank)
    add_title(slide, "평가 참여 구조", "1조·2조 외에 가상 비교군 0조·4조를 포함해 비교 평가의 직관성을 높였다")
    add_text(slide, 1.35, 3.05, 29.8, 1.0, "3조는 조장 협의를 통해 온톨로지 & RAG 벤치마크 평가에서는 제외하였고, 별도 테크니컬 인사이트를 도출할 예정이다.", 17, True, COLORS["navy"])
    add_rule(slide, 1.35, 4.55, 29.5, COLORS["line"])
    add_text(slide, 1.45, 5.35, 10.0, 0.6, "평가 참여", 13, True, COLORS["blue"])
    rows = [
        ["구분", "조", "구성/해석"],
        ["실제 참여", "1조", "온톨로지. 사실상 온톨로지+RAG를 함께 사용한 하이브리드 방식"],
        ["실제 참여", "2조", "Advanced RAG"],
        ["가상 비교군", "0조", "순수 RAG"],
        ["가상 비교군", "4조", "하이브리드 RAG. 온톨로지+전통 RAG"],
    ]
    add_simple_table(slide, 1.45, 6.25, [4.0, 3.0, 20.5], rows, True, 9.4)
    add_text(slide, 2.0, 13.95, 28.5, 1.0, "가상 조는 특정 조의 우열을 강조하기보다, 순수 RAG와 하이브리드 RAG의 기준선을 함께 놓기 위한 비교 장치다.", 18, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3 transition
    slide = prs.slides.add_slide(blank)
    add_title(slide, "평가 체계가 바뀌었다", "정답 보정과 함정 문항 처리로 단순 점수 비교에서 서비스 품질 평가로 이동")
    rows = [
        ["구분", "답변 보정전 산정", "답변 보정후 산정"],
        ["입력", "AI Lab API 정확도 산정(답변 보정전)", "AI Lab API 정확도 산정(답변 보정후)"],
        ["정답", "초기 예상 답변", "예상 답변·근거·핵심포함어 보정"],
        ["STD-S", "Snowflake RAG 문항으로 채점", "문서에 없는 함정 질문으로 처리"],
        ["평가 의미", "답변 유사도 중심", "근거 부재 감지 + 답변 정확도"],
    ]
    add_simple_table(slide, 1.2, 3.0, [4.2, 12.3, 12.3], rows, True, 9.5)
    add_change_chart(slide, 1.5, 8.35, 15.0, 7.3, data["old_scores"], data["new_scores"])
    add_bullets(slide, 18.2, 8.6, 12.0, 5.7, [
        "기존 0~2조 답변 원문은 보존하고, 답변 보정후 기준으로 점수만 재산정했다.",
        "4조는 ont_platform v4 답변을 추가해 동일 기준으로 비교했다.",
        "정답 자체의 편향과 질문 유도 오류를 보정하면서 평가 목표가 더 현실적인 서비스 검증으로 바뀌었다.",
    ], 12)
    add_footer(slide, 3)

    # 4 ranking
    slide = prs.slides.add_slide(blank)
    add_title(slide, "최종 순위는 4조, 1조, 2조, 0조", "1조 wrong-source 감점과 2조 STD-S-05 보정을 모두 반영한 최종 점수")
    add_team_chart(slide, 1.3, 3.0, 17.0, 10.7, data["new_scores"])
    rows = [
        ["팀", "최종 정확도", "해석"],
        ["4조", "85.00%", "ont_platform v4 기준, 함정 질문 거절과 출처 제한에서 가장 안정적"],
        ["1조", "76.88%", "RAG 하이브리드 효과는 있으나, STD-S wrong-source 7문항 감점 반영"],
        ["2조", "72.92%", "Advanced RAG 문항은 강하고 STD-S-05는 적절했으나, 다른 STD-S 차단 한계 존재"],
        ["0조", "70.42%", "순수 RAG로 기준선 역할. 보정 후 도입 고려 수준"],
    ]
    add_simple_table(slide, 19.0, 3.0, [2.6, 2.7, 8.0], rows, True, 8.4)
    add_footer(slide, 3)

    # 4 category heatmap
    slide = prs.slides.add_slide(blank)
    add_title(slide, "카테고리별 강점은 다르게 나타났다", "Advanced RAG는 2조, Ontology와 Snowflake 함정 대응은 1조·4조가 상대적으로 강하다")
    add_heatmap(slide, 1.4, 3.0, data["categories"])
    add_text(slide, 1.4, 7.8, 14.8, 0.55, "카테고리 평균", 11, True, COLORS["muted"])
    add_simple_table(slide, 1.4, 8.55, [4.3, 3.0, 3.0, 3.0, 3.0, 3.0], [
        ["카테고리", "0조", "1조", "2조", "4조", "평균"],
        *[[c["category"], f'{c["Team0"]:.2f}', f'{c["Team1"]:.2f}', f'{c["Team2"]:.2f}', f'{c["Team4"]:.2f}', f'{c["avg"]:.2f}'] for c in data["categories"]],
    ], True, 8.2)
    add_bullets(slide, 20.5, 3.2, 10.0, 8.0, [
        "2조는 Advanced RAG 특화 질문에서 93.75%로 가장 강했다.",
        "1조는 전 영역에서 균형적으로 높았지만 wrong-source 보정 후 순위가 낮아졌다.",
        "4조는 Snowflake 함정 문항에서 90.62%로 강하지만 Ontology 문항은 추가 개선 여지가 있다.",
    ], 12.2)
    add_footer(slide, 4)

    # 5 team diagnosis
    slide = prs.slides.add_slide(blank)
    add_title(slide, "팀별 해석은 ‘기술명’보다 실제 동작 기준으로 봐야 한다", "평가 명칭, 구현 실제, 출처 차단 여부가 모두 다르므로 점수만으로 결론을 내리면 위험하다")
    rows = [
        ["팀", "설계/구현 성격", "강점", "한계"],
        ["0조", "순수 RAG", "기본 검색·생성 기준선", "관련 없음 감지와 정교한 근거 통제 한계"],
        ["1조", "온톨로지 취지 + RAG 하이브리드", "하이브리드 효과와 균형형 성능", "HNIX 중간점검 PPT 등 잘못된 출처를 근거로 사용"],
        ["2조", "Advanced RAG", "Advanced RAG 문항 최상위, STD-S-05는 적절", "다수 STD-S에서 다른 문서/일반론 출력. threshold 또는 환각 가능"],
        ["4조", "ont_platform v4", "함정 질문 대응 강함", "Ontology 문항의 개념·관계 설명 보강 필요"],
    ]
    add_simple_table(slide, 1.2, 3.0, [2.6, 7.0, 7.4, 11.4], rows, True, 8.6)
    add_text(slide, 1.4, 13.2, 29.0, 1.1, "진짜 비교 포인트는 ‘어떤 시스템이 평균 점수가 높은가’보다 ‘어떤 질문 유형에서 어떤 검색 기술이 실패했는가’다.", 19, True, COLORS["navy"], PP_ALIGN.CENTER)
    add_footer(slide, 5)

    # 6 STD-S example
    slide = prs.slides.add_slide(blank)
    add_title(slide, "STD-S 문항은 Snowflake 성능 문제가 아니라 함정 질문 처리 문제였다", "논문에 Snowflake 구현 내용이 없으므로, 정답은 관련 없음/문서 근거 없음의 명시다")
    add_text(slide, 1.35, 3.0, 28.5, 0.85, "STD-S01 질문 예시", 12, True, COLORS["blue"])
    add_text(slide, 1.35, 3.75, 29.8, 1.15, "Snowflake 기반 RAG 평가에서 문서 저장소와 근거 페이지를 왜 함께 관리해야 하는가?", 21, True, COLORS["ink"])
    add_rule(slide, 1.35, 5.45, 30.0, COLORS["line"])
    add_text(slide, 1.4, 6.15, 8.5, 0.55, "만점 답변의 핵심", 11, True, COLORS["green"])
    add_bullets(slide, 1.4, 6.9, 9.0, 5.2, [
        "Snowflake RAG의 구체 구현은 제공 문서에 없다.",
        "따라서 관련 없음 또는 문서 근거 없음이라고 명시한다.",
        "부가 설명이 없어도 100점, 일반적 RAG 관점 설명을 덧붙여도 100점이다.",
    ], 11.4)
    add_text(slide, 12.2, 6.15, 8.5, 0.55, "감점되는 답변", 11, True, COLORS["orange"])
    add_bullets(slide, 12.2, 6.9, 9.0, 5.2, [
        "관련 없음 표현 없이 Snowflake 기능을 실제 근거처럼 설명한다.",
        "Snowflake 언급 질문인데 문서 부재를 말하지 않는다.",
        "온톨로지/RAG 등 다른 주제로 우회 답변한다.",
    ], 11.4)
    add_text(slide, 22.5, 6.15, 8.5, 0.55, "실서비스 정책", 11, True, COLORS["navy"])
    add_bullets(slide, 22.5, 6.9, 9.0, 5.2, [
        "지원 근거 없음 감지",
        "부분 근거가 있으면 범위 제한 설명",
        "근거 없는 일반지식 생성은 옵션화 또는 차단",
    ], 11.4)
    add_footer(slide, 6)

    # 7 wrong-source finding
    slide = prs.slides.add_slide(blank)
    add_title(slide, "추가 발견: wrong-source leakage가 있었다", "STD-S 함정 질문에서 ‘관련 없음’만 보는 평가는 부족하다. 잘못된 출처를 끌어온 답변도 감점해야 한다")
    add_text(slide, 1.35, 3.0, 28.5, 0.7, "관찰된 현상", 12, True, COLORS["blue"])
    rows = [
        ["대상", "관찰", "평가상 의미"],
        ["1조", "Snowflake 직접 근거는 없다고 말한 뒤 HNIX AI Labs 2기 중간점검 자료를 근거처럼 사용", "관련 없음 표현만으로 만점 처리하면 점수 과대 산정"],
        ["2조", "STD-S-05는 문서 내 명확한 정보가 없다고 제한해 적절했으나, 다른 STD-S 일부에서 다른 논문/일반론/RDF/온톨로지 내용을 출력", "문서 범위 차단 실패 또는 환각 가능성"],
        ["평가자", "wrong-source가 포함되어도 관련 없음 문구가 있으면 100점을 부여", "judge/evaluator도 출처 검증 로직이 필요"],
    ]
    add_simple_table(slide, 1.35, 4.1, [3.0, 15.0, 10.5], rows, True, 8.4)
    add_text(slide, 2.2, 10.5, 27.0, 1.0, "함정 질문의 만점 조건은 ‘관련 없음 명시’뿐 아니라 ‘허용된 평가 문서 밖의 근거를 사용하지 않음’까지 포함해야 한다.", 19, True, COLORS["red"], PP_ALIGN.CENTER)
    add_bullets(slide, 3.0, 12.4, 25.5, 3.5, [
        "출처 문서 whitelist를 두고, 검색·생성·평가 단계에서 모두 검증한다.",
        "답변이 문서 부재를 말하더라도 잘못된 문서 내용을 근거로 확장하면 감점한다.",
        "1조는 STD-S wrong-source 7개 문항을 100→60으로 감점하여 최종 76.88%로 반영했다.",
    ], 12.4)
    add_footer(slide, 7)

    # 8 implication 1
    slide = prs.slides.add_slide(blank)
    add_title(slide, "시사점 1: ‘관련 없음’과 ‘부분 설명 가능’을 분리해야 한다", "잘못된 질문이 들어와도 시스템은 무조건 답하지 말고 근거 범위를 먼저 판단해야 한다")
    add_text(slide, 1.4, 3.1, 11.5, 0.6, "평가/서비스 응답 정책", 12, True, COLORS["muted"])
    rows = [
        ["질문 상태", "응답 원칙", "점수/서비스 의미"],
        ["문서에 전혀 없음", "관련 없음·근거 없음 명시", "정답 또는 안전 응답"],
        ["질문 일부는 없음", "없는 부분은 제외하고 문서 내 근거만 설명", "부분 정답 + 범위 제한"],
        ["문서에 있음", "문서·페이지·청크 근거와 함께 설명", "정상 답변"],
        ["일반 지식 필요", "옵션으로 일반 지식/전문가 관점 분리", "고급 서비스 모드"],
    ]
    add_simple_table(slide, 1.4, 4.05, [5.0, 11.0, 10.4], rows, True, 9)
    add_text(slide, 3.5, 13.4, 25.0, 1.0, "핵심은 거절이 아니라 근거 범위의 정직한 표시다.", 22, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide, 8)

    # 9 implication 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "시사점 2: 예상 답변도 편향될 수 있다", "고객이 원하는 것은 문서 복원이 아니라, 문서 기반 답변과 전문가적 보정의 균형일 수 있다")
    add_text(slide, 1.4, 3.0, 29.0, 1.2, "전문가는 때로 ‘우문현답’을 한다. RAG 서비스도 문서 기반 답변만 고집할지, 문서의 한계를 밝히고 보편적 통용 지식을 덧붙일지 옵션을 가져야 한다.", 18, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_rule(slide, 5.2, 5.35, 23.0, COLORS["line"])
    rows = [
        ["모드", "답변 방식", "사용 장면"],
        ["Strict Evidence", "제공 문서에 있는 내용만 답변", "시험/감사/컴플라이언스"],
        ["Evidence + Caveat", "문서 기반 답변 후 주관성·한계 표시", "의사결정 보고"],
        ["Expert Expansion", "문서 근거와 일반 통용 지식 분리 제시", "컨설팅/고급 분석 서비스"],
    ]
    add_simple_table(slide, 3.0, 7.0, [5.2, 12.0, 10.0], rows, True, 9.2)
    add_footer(slide, 9)

    # 10 implication 3+4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "시사점 3·4: 정확도는 답변만으로 판단하기 어렵다", "여러 논문·여러 견해·여러 검색 후보가 있을수록 검색 품질과 생성 품질을 분리해서 봐야 한다")
    add_text(slide, 1.5, 3.1, 9.0, 0.5, "평가 한계", 11, True, COLORS["orange"])
    add_bullets(slide, 1.5, 3.85, 10.5, 6.4, [
        "정형화된 문제는 답이 하나처럼 보일 수 있다.",
        "여러 논문에서 질문하면 견해 차이와 비슷한 비중의 후보가 동시에 검색된다.",
        "이때 단일 예상 답변과 비교하는 방식은 정확도 평가에 한계가 있다.",
    ], 12)
    add_text(slide, 14.0, 3.1, 9.0, 0.5, "분리 측정", 11, True, COLORS["blue"])
    add_bullets(slide, 14.0, 3.85, 10.5, 6.4, [
        "검색이 올바른 문서/페이지/청크를 찾았는가",
        "생성이 검색 근거를 충실히 반영했는가",
        "근거 없는 일반화나 환각이 있었는가",
    ], 12)
    add_text(slide, 1.8, 12.8, 29.0, 1.0, "따라서 평가 결과에는 답변 점수뿐 아니라 문서명, 페이지, 청크 ID, 검색 점수, 사용 여부가 함께 남아야 한다.", 18, True, COLORS["navy"], PP_ALIGN.CENTER)
    add_footer(slide, 10)

    # 11 implication 5+6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "시사점 5·6: 하이브리드 검색은 관측 가능해야 개선된다", "벡터DB, BM25, 온톨로지 검색의 기여를 분리하지 않으면 성능 향상 지점을 찾기 어렵다")
    rows = [
        ["필요 로그", "왜 필요한가", "개선 포인트"],
        ["검색 기술", "어떤 엔진이 근거를 찾았는지 식별", "벡터/BM25/온톨로지별 튜닝"],
        ["검색 점수", "threshold 문제와 노이즈 유입 확인", "유사도 임계값 조정"],
        ["문서·페이지·청크", "답변 근거의 재현성 확보", "근거 적합성 평가"],
        ["생성 사용 여부", "검색됐지만 사용되지 않은 근거 파악", "프롬프트/컨텍스트 구성 개선"],
        ["엔진 선택 옵션", "개별 성능과 조합 성능 분리 측정", "디버깅과 A/B 테스트"],
    ]
    add_simple_table(slide, 1.3, 3.0, [5.0, 10.3, 10.5], rows, True, 8.8)
    add_text(slide, 3.0, 14.15, 27.0, 0.9, "이상적인 하이브리드 답변은 각 검색 기술의 개별 성능이 충분히 높을 때 만들어진다.", 19, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12 operating model
    slide = prs.slides.add_slide(blank)
    add_title(slide, "향후 평가 프레임워크", "정답 보정, 출처 whitelist, 근거 추적, 검색기술별 로그, 응답 모드 분리를 표준으로 가져가야 한다")
    rows = [
        ["단계", "평가 항목", "산출물"],
        ["1. 질문 분류", "문서 내/문서 외/부분 근거/일반 지식 필요", "question_type, support_scope"],
        ["2. 출처 검증", "허용 문서 whitelist, wrong-source leakage", "source_validity_flag"],
        ["3. 검색 평가", "문서·페이지·청크 recall, 검색기술별 hit", "retrieval_trace"],
        ["4. 생성 평가", "근거 반영률, 환각, 범위 제한 표현", "answer_score, hallucination_flag"],
        ["5. 하이브리드 분석", "벡터/BM25/온톨로지별 기여도", "engine_contribution"],
        ["6. 운영 개선", "threshold, reranker, prompt, ontology mapping", "improvement backlog"],
    ]
    add_simple_table(slide, 1.4, 3.0, [4.2, 12.0, 12.0], rows, True, 9)
    add_bullets(slide, 2.0, 12.4, 27.0, 2.8, [
        "답변 보정후 산정은 점수 보정이 아니라 평가 철학의 보정이다.",
        "서비스화 단계에서는 ‘정답을 말하는 능력’과 ‘모르면 모른다고 말하는 능력’을 같은 수준으로 평가해야 한다.",
    ], 13)
    add_footer(slide, 12)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build_deck())
