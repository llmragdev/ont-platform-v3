#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
기술이전 프레젠테이션 PowerPoint 생성 스크립트

사용:
    python scripts/create_presentation.py

생성 파일:
    ont_platform/v3/TECHNOLOGY_TRANSFER_PRESENTATION.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import sys

# 기본 설정
TITLE_COLOR = RGBColor(0, 51, 102)  # 진한 파란색
ACCENT_COLOR = RGBColor(0, 102, 204)  # 밝은 파란색
TEXT_COLOR = RGBColor(51, 51, 51)  # 어두운 회색

def add_title_slide(prs, title, subtitle=""):
    """제목 슬라이드 추가"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    if subtitle:
        subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_list):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for item in content_list:
        if isinstance(item, str):
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
        elif isinstance(item, tuple):
            # (text, level) 형태
            p = tf.add_paragraph()
            p.text = item[0]
            p.level = item[1] if len(item) > 1 else 0
            p.font.size = Pt(16) if item[1] == 1 else Pt(18)
            p.font.color.rgb = TEXT_COLOR

    return slide

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """2열 슬라이드"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    # 좌측 컬럼
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4.25), Inches(5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True

    p = left_tf.add_paragraph()
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    for item in left_content:
        p = left_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(6)

    # 우측 컬럼
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(4.5), Inches(5))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True

    p = right_tf.add_paragraph()
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    for item in right_content:
        p = right_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(6)

    return slide

def create_presentation():
    """프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("[*] PowerPoint 프레젠테이션 생성 중...")

    # Slide 1: Title
    add_title_slide(prs,
        "운영형 온톨로지 플랫폼",
        "조선·제조 산업용 지능형 데이터 관리 및 의사결정 시스템\n\n기술이전 설명서 | 2026년 5월")
    print("[OK] Slide 1: 제목")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary (30초 요약)", [
        ("문제", 1),
        ("조선/제조 데이터의 복잡성 증가", 2),
        ("기존 솔루션 성능/비용 문제", 2),
        ("", 0),
        ("해결책", 1),
        ("PostgreSQL 기반 운영 온톨로지 저장소", 2),
        ("W3C SPARQL 표준 호환 쿼리 계층", 2),
        ("업무 액션 & write-back 기능", 2),
        ("", 0),
        ("가치", 1),
        ("성능 개선: 200배 (20초 → 100ms)", 2),
        ("개발 기간: 4주", 2),
        ("인프라 비용: 무료~저비용", 2),
    ])
    print("[OK] Slide 2: Executive Summary")

    # Slide 3: 시장 현황
    add_content_slide(prs, "시장 현황 & 경쟁 분석", [
        ("Palantier Foundry", 1),
        ("강점: 엔터프라이즈급 성능 | 약점: 고비용($1M+), 폐쇄 시스템", 2),
        ("RDF 준수: X (자체 Object-Link-Action 모델)", 2),
        ("", 0),
        ("Neo4j + 그래프DB", 1),
        ("강점: 오픈소스, 그래프 성능 | 약점: Cypher 강제, SPARQL 미지원", 2),
        ("RDF 준수: ? (플러그인 의존)", 2),
        ("", 0),
        ("국내 경쟁 솔루션", 1),
        ("솔트룩스(온톨로지+LLM), SKAI(GraphRAG), NUBISON(산업 AX)", 2),
    ])
    print("[OK] Slide 3: 시장 현황")

    # Slide 4: 기술 진단
    add_content_slide(prs, "RDF 성능 한계 분석", [
        ("구조적 문제 (구현 개선 불가)", 1),
        ("", 0),
        ("JSONL + rdflib:", 1),
        ("10K: O 100ms | 100K: X 20-60초 | 1M: X 분 단위", 2),
        ("원인: 조인 폭발, 메모리 스캔", 2),
        ("", 0),
        ("PostgreSQL + SQL (제안):", 1),
        ("10K: O 10ms | 100K: O 100ms | 1M: O < 1초", 2),
        ("이유: 인덱스 + SQL 최적화", 2),
        ("", 0),
        ("해결: 저장소 계층 전환 필수", 1),
    ])
    print("[OK] Slide 4: 기술 진단")

    # Slide 5: 솔루션 아키텍처
    add_content_slide(prs, "하이브리드 3계층 아키텍처", [
        ("Layer 1: 표준 API 계층", 1),
        ("W3C SPARQL 1.1 (read-only), JSON-LD/Turtle/N-Triples export", 2),
        ("", 0),
        ("Layer 2: 번역 & 엔진 계층", 1),
        ("SPARQL→SQL 번역기 (hot-path), rdflib fallback", 2),
        ("Action/Write-back 엔진", 2),
        ("", 0),
        ("Layer 3: 저장소 계층", 1),
        ("PostgreSQL: Operational Ontology Store", 2),
        ("Optional: Neo4j (graph acceleration)", 2),
    ])
    print("[OK] Slide 5: 솔루션 아키텍처")

    # Slide 6: 기술 차별화
    add_content_slide(prs, "5가지 차별화 포인트", [
        ("1. 운영형 온톨로지 (Operational Ontology)", 1),
        ("읽기만 가능 X → 업무 상태 변경 가능 O", 2),
        ("BOM 변경, 도면 개정, 검사 결과 반영 자동화", 2),
        ("", 0),
        ("2. Write-back & Action", 1),
        ("조회만 가능 X → 조회+실행+이력 추적 O", 2),
        ("", 0),
        ("3. 도메인 특화 스키마", 1),
        ("범용 RDF X → 조선/제조 네이티브 모델 O", 2),
        ("", 0),
        ("4. 감사 추적 & 혈통", 1),
        ("누가, 무엇을, 왜 변경했는지 추적", 2),
        ("", 0),
        ("5. 표준 호환성 (선택적)", 1),
        ("DBpedia, Wikidata 등 외부 온톨로지 통합 가능", 2),
    ])
    print("[OK] Slide 6: 차별화 포인트")

    # Slide 7: 성능 목표
    add_content_slide(prs, "Hot-Path 쿼리 성능 목표", [
        ("100K entities, 1M relationships 기준", 1),
        ("", 0),
        ("Simple Lookup: < 50ms O", 0),
        ("예: 선박 ID, 부품 시리얼 번호 조회", 1),
        ("", 0),
        ("Indexed Filter: < 200ms O", 0),
        ("예: 상태별, 타입별 엔티티 조회", 1),
        ("", 0),
        ("One-Hop Relation: < 300ms O", 0),
        ("예: 선박의 모든 블록, 블록의 모든 부품", 1),
        ("", 0),
        ("Two-Hop Relation: < 1s O", 0),
        ("예: 부품의 공급자, 공급자의 다른 부품", 1),
        ("", 0),
        ("복잡 쿼리: Async/Batch 처리", 0),
        ("예: 영향 분석, RDF export", 1),
    ])
    print("[OK] Slide 7: 성능 목표")

    # Slide 8: 4주 로드맵
    add_content_slide(prs, "4주 구현 로드맵", [
        ("Week 1 (05-27~31): 기초 구축", 1),
        ("rdflib SPARQL 통합, PostgreSQL 스키마, 개발 환경", 2),
        ("산출물: 30개 SPARQL 테스트, DDL, Docker", 2),
        ("", 0),
        ("Week 2 (06-03~07): 번역 엔진", 1),
        ("SPARQL→SQL 번역기 (500줄), 50개 패턴", 2),
        ("산출물: 500줄 코드, 50개 E2E 테스트", 2),
        ("", 0),
        ("Week 3 (06-10~14): API 통합", 1),
        ("FastAPI 엔드포인트, 트랜잭션, Changelog", 2),
        ("산출물: 5개 API, 10개 동시성 테스트", 2),
        ("", 0),
        ("Week 4 (06-17~21): 성능 검증", 1),
        ("100K-1M 벤치마크, 최적화, 최종 문서", 2),
        ("산출물: 성능 보고서, 운영 가이드", 2),
    ])
    print("[OK] Slide 8: 4주 로드맵")

    # Slide 9: 기술 이전 옵션
    add_content_slide(prs, "기술이전 3가지 옵션", [
        ("Option 1: 기술 라이선싱 (기술이전)", 1),
        ("가격: $200K-500K | 기간: 4주 구현 + 2주 이전", 2),
        ("포함: 소스코드, 문서, 테스트, 컨설팅", 2),
        ("", 0),
        ("Option 2: 제품화 협력 (50:50 수익 배분)", 1),
        ("기간: 6개월 | 역할: 기술+마케팅 협력", 2),
        ("수익: 초기 계약금 + 라이선싱료 배분", 2),
        ("", 0),
        ("Option 3: SaaS 운영 (호스팅 서비스)", 1),
        ("가격: Base $10K/월 + 추가 $100/1M entities", 2),
        ("대상: 중소 조선사, 협력사", 2),
    ])
    print("[OK] Slide 9: 기술이전 옵션")

    # Slide 10: 경제성 분석
    add_content_slide(prs, "경제성 분석 (5년 기준)", [
        ("vs Palantier Foundry", 1),
        ("Palantier: $7.5M+ | 제안: $1M-1.5M | 절감: 80-85% O", 2),
        ("", 0),
        ("vs Neo4j + SI 구축", 1),
        ("Neo4j+SI: $800K-1M | 제안: $550K | 절감: 30-45% O", 2),
        ("", 0),
        ("ROI (기술이전 회사 관점)", 1),
        ("개발비: $112K | 첫 고객 계약금: $300K", 2),
        ("회수 기간: 3개월 내 | 연간 순익: $200K+", 2),
    ])
    print("[OK] Slide 10: 경제성 분석")

    # Slide 11: 경쟁 우위
    add_content_slide(prs, "경쟁 제품 비교", [
        ("가격", 1),
        ("Palantier: $1M+ vs 제안: $300K (85% 저가)", 2),
        ("", 0),
        ("RDF 표준 준수", 1),
        ("제안만 유일하게 W3C SPARQL 호환 O", 2),
        ("", 0),
        ("Write-back & 감사 추적", 1),
        ("경쟁사는 읽기 중심 | 제안은 실행+추적 O", 2),
        ("", 0),
        ("도메인 특화", 1),
        ("조선/제조 네이티브 스키마 (경쟁사 없음) O", 2),
        ("", 0),
        ("구현 기간", 1),
        ("경쟁사: 3-6개월 vs 제안: 4주 O", 2),
    ])
    print("[OK] Slide 11: 경쟁 우위")

    # Slide 12: 포지셔닝
    add_content_slide(prs, "시장 포지셔닝", [
        ("피해야 할 경쟁지", 1),
        ("X 범용 온톨로지 (솔트룩스 직경쟁)", 2),
        ("X GraphRAG 검색 (SKAI와 경쟁)", 2),
        ("X Palantier 성능 모방", 2),
        ("", 0),
        ("승리 지점", 1),
        ("O 조선/제조 산업 특화 (도메인 네이티브)", 2),
        ("O 운영 온톨로지 + Write-back (실행 가능)", 2),
        ("O 감사 추적 + Lineage (컴플라이언스)", 2),
        ("O 표준 호환 + 경제성 (SPARQL + 저가)", 2),
    ])
    print("[OK] Slide 12: 포지셔닝")

    # Slide 13: 고객 시나리오
    add_content_slide(prs, "고객 가치 시나리오", [
        ("대형 조선사", 1),
        ("현황: 변경 영향 분석 2-3주 수동 작업", 2),
        ("해결: 자동 영향 분석 1시간, 감사 추적 기록", 2),
        ("ROI: 연 $500K (생산성 + 컴플라이언스)", 2),
        ("", 0),
        ("중소 협력사", 1),
        ("현황: Excel 기반 관리, AI 정확도 낮음", 2),
        ("해결: 경량 온톨로지 SaaS, 정확한 데이터 기반 AI", 2),
        ("비용: $10K/월 | ROI: 연 $150K+ (위험 감지)", 2),
    ])
    print("[OK] Slide 13: 고객 시나리오")

    # Slide 14: 기술 스펙
    add_content_slide(prs, "핵심 기술 스펙", [
        ("Language & Framework", 1),
        ("Python 3.11+, FastAPI, SQLAlchemy, rdflib", 2),
        ("", 0),
        ("Database", 1),
        ("PostgreSQL 14+ (primary), Optional Neo4j", 2),
        ("", 0),
        ("Standards", 1),
        ("W3C SPARQL 1.1, RDF 1.1, SHACL, JSON-LD", 2),
        ("", 0),
        ("Scalability", 1),
        ("100K-10M entities, 1M-100M relationships", 2),
        ("100-1000 QPS, 10-100 concurrent users", 2),
        ("", 0),
        ("Deployment", 1),
        ("Docker, Kubernetes, Cloud (AWS/GCP/Azure), On-premise", 2),
    ])
    print("[OK] Slide 14: 기술 스펙")

    # Slide 15: 다음 단계
    add_content_slide(prs, "기술이전 다음 단계", [
        ("Step 1: 기술 수용 의사 확인", 1),
        ("Step 2: 라이선싱 모델 선택 (Option 1/2/3)", 1),
        ("Step 3: 계약 및 NDA 논의", 1),
        ("Step 4: 파일럿 고객 확보 (2026년 8-9월)", 1),
        ("Step 5: 구현 및 수익화", 1),
        ("", 0),
        ("기술이전 일정", 1),
        ("5월: 기술 완성 | 7월: 이전 | 8-9월: 파일럿 | 10월: 시장 진출", 2),
    ])
    print("[OK] Slide 15: 다음 단계")

    # Slide 16: 핵심 메시지
    add_content_slide(prs, "핵심 메시지", [
        ("조선/제조 데이터를 안전하게 이해하고", 1),
        ("실행하는 플랫폼", 1),
        ("", 0),
        ("O 성능: Palantier 급 (200배 개선)", 0),
        ("O 표준: W3C SPARQL 호환", 0),
        ("O 비용: 85% 절감", 0),
        ("O 속도: 4주 구현", 0),
        ("", 0),
        ("메시지 전환", 1),
        ("부정확한 AI 답변 감소 (약함)", 2),
        ("→ 업무 상태 변경+추적 플랫폼 (강함)", 2),
    ])
    print("[OK] Slide 16: 핵심 메시지")

    # Slide 17: Contact
    add_content_slide(prs, "기술이전 문의", [
        ("기술 담당자", 1),
        ("Architecture: Claude Code", 2),
        ("기술 검토: Antigravity", 2),
        ("도메인 전문성: Kodex", 2),
        ("", 0),
        ("제공 자료", 1),
        ("문서: 5,200줄 (분석, 설계, 스펙)", 2),
        ("코드: 1,500줄 (Python, SQL, Bash)", 2),
        ("테스트: 150개+ (단위, 통합, 성능)", 2),
        ("", 0),
        ("기술이전 일정", 1),
        ("[DATE] 2026년 5월 기술 완성", 2),
        ("[DATE] 2026년 7월 이전 & 교육", 2),
        ("[DATE] 2026년 8-9월 파일럿 고객", 2),
    ])
    print("[OK] Slide 17: Contact")

    # 저장
    output_path = "E:\\ontology_edu\\X_ont_std\\ont_platform\\v3\\TECHNOLOGY_TRANSFER_PRESENTATION.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    print(f"\n[OK] 완료!")
    print(f"[FILE] {output_path}")
    print(f"[SLIDES] {len(prs.slides)}개")

if __name__ == '__main__':
    try:
        create_presentation()
    except Exception as e:
        print(f"[ERROR] {e}", file=sys.stderr)
        sys.exit(1)
